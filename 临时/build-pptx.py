from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

ASSETS = "assets"

def add_full_slide(slide, img_path):
    """Background image covers entire slide."""
    slide.shapes.add_picture(img_path, 0, 0, W, H)

def add_split_slide(slide, img_path, img_on_right, title, bullets):
    """Image on one side (40%), text on the other."""
    img_width = int(W * 0.4)
    img_left = W - img_width if img_on_right else 0
    # 40% width, full height
    slide.shapes.add_picture(img_path, img_left, 0, img_width, H)

    # Text box on the other side (60% wide)
    text_left = 0 if img_on_right else img_width
    text_width = W - img_width
    txBox = slide.shapes.add_textbox(text_left, Inches(1.0), text_width, H - Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)
    p.alignment = PP_ALIGN.CENTER

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

def add_overlay_slide(slide, img_path, title, text_lines):
    """Full bg image with white overlay text at center."""
    slide.shapes.add_picture(img_path, 0, 0, W, H)

    # Semi-transparent dark overlay for readability
    overlay = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(2), Inches(2.5), Inches(9.333), Inches(2.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.fill.fore_color.brightness = 0.0
    overlay.line.fill.background()
    # transparency via alpha isn't directly supported, use a dark shape

    txBox = slide.shapes.add_textbox(Inches(2.5), Inches(2.6), Inches(8.333), Inches(2.4))
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    for line in text_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

# Slide 1: Cover (full image, no text)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
add_full_slide(slide, f"{ASSETS}/slide-01.png")

# Slide 2: Once upon a time
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-02.jpg", True,
    "Once upon a time ...",
    ["There was a little frog.",
     "He lived in a well.",
     "He was very happy."])

# Slide 3: The Frog's World
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-03.jpg", False,
    "The Frog's World",
    ['The frog looked up.',
     'He saw a small blue sky.',
     'He said: "The world is small!',
     'I am the king of the world!"'])

# Slide 4: A Bird Comes
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-04.jpg", True,
    "A Bird Comes",
    ["One day, a bird flew to the well.",
     "The bird said:",
     '"Hello, little frog!"',
     '"Come out and see the big world!"'])

# Slide 5: The Frog is Surprised
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-05.png", False,
    "The Frog is Surprised",
    ['The frog said:',
     '"No, no! The world is small!"',
     '"I can see it all from here!"'])

# Slide 6: The Bird Tells Him
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-06.png", True,
    "The Bird Tells Him",
    ['"Come and see!"',
     "There are big seas, tall trees,",
     "and blue skies far away.",
     "The world is very, very big!"])

# Slide 7: The Frog Learns
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-07.png", False,
    "The Frog Learns",
    ["The frog jumped out of the well.",
     "He looked around.",
     "He saw the big blue sky.",
     "He saw tall green trees.",
     "He saw a beautiful lake."])

# Slide 8: The Frog Understands
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_split_slide(slide, f"{ASSETS}/slide-08.png", True,
    "The Frog Understands",
    ['"Oh! The world is so big!"',
     "I only saw a small part of it.",
     "I was wrong."])

# Slide 9: The Lesson
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_overlay_slide(slide, f"{ASSETS}/slide-09.png",
    "The Lesson – 井底之蛙",
    ["A frog in a well cannot see the big ocean.",
     "A small view gives a small mind."])

# Slide 10: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_overlay_slide(slide, f"{ASSETS}/slide-10.png",
    "Thank You!",
    [])

output = "井底之蛙.pptx"
prs.save(output)
print(f"Saved: {output} ({len(prs.slides)} slides)")
